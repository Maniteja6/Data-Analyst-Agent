"""PPTXGenerator — exports InsightReport to a branded PowerPoint deck.

Real-time design:
    PPTX generation runs in a thread pool executor and emits
    ``report:slide_complete`` Socket.IO events as each slide is built.

Deck structure (7 slides):
    Slide 1: Title — dataset name, date, subtitle
    Slide 2: Executive Summary
    Slide 3: Data Quality Overview
    Slide 4-5: Top Insights (2-3 per slide)
    Slide 6: Recommendations
    Slide 7: Next Steps + Q&A

Design tokens:
    Primary: #5B4FE8 (DataPilot violet)
    Dark:    #1A1A3E (DataPilot navy)
    Light:   #F5F4F1 (DataPilot canvas)
"""

from __future__ import annotations

import asyncio
import contextlib
import io
from datetime import datetime
from typing import TYPE_CHECKING, Any

import structlog

if TYPE_CHECKING:
    from pptx.shapes.autoshape import Shape
    from pptx.slide import Slide

logger = structlog.get_logger(__name__)


async def export_to_pptx(
    report: dict[str, Any],
    dataset_name: str = "Dataset",
    sio: Any = None,  # noqa: ANN401
    dataset_id: str = "",
) -> bytes:
    """Generate a PPTX deck from an InsightReport dict.

    Returns:
        Raw PPTX bytes ready for S3 upload.
    """
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, _build_deck, report, dataset_name, sio, dataset_id, loop
    )


def _build_deck(
    report: dict[str, Any],
    dataset_name: str,
    sio: Any,  # noqa: ANN401
    dataset_id: str,
    loop: asyncio.AbstractEventLoop,
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError:
        return b"[PPTX export requires python-pptx: pip install python-pptx]"

    slide_done = [0]

    def _emit(slide_num: int, label: str) -> None:
        slide_done[0] += 1
        if sio and dataset_id:
            with contextlib.suppress(Exception):
                asyncio.run_coroutine_threadsafe(
                    sio.emit(
                        "report:slide_complete",
                        {
                            "dataset_id": dataset_id,
                            "slide": slide_num,
                            "label": label,
                        },
                        room=f"dataset:{dataset_id}",
                    ),
                    loop,
                )

    violet = RGBColor(0x5B, 0x4F, 0xE8)
    navy = RGBColor(0x1A, 0x1A, 0x3E)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    def _text_box(
        slide: Slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int = 18,
        bold: bool = False,
        color: RGBColor = None,
    ) -> Shape:
        tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = bold
        if color:
            p.runs[0].font.color.rgb = color
        return tx_box

    def _fill_bg(slide: Slide, color_hex: str = "1A1A3E") -> None:
        """Fill slide background with a solid colour."""
        from pptx.dml.color import RGBColor

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex)

    # ── Slide 1: Title ────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    _fill_bg(s1, "1A1A3E")
    _text_box(s1, "DataPilot", 0.5, 0.5, 6, 1, size=12, color=white)
    _text_box(s1, dataset_name, 0.5, 1.2, 12, 2, size=36, bold=True, color=white)
    _text_box(s1, "AI-Powered Data Analysis Report", 0.5, 3.2, 10, 1, size=18, color=violet)
    _text_box(s1, datetime.utcnow().strftime("%B %Y"), 0.5, 4.2, 6, 0.5, size=14, color=white)
    _emit(1, "Title")

    # ── Slide 2: Executive Summary ────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _text_box(s2, "Executive Summary", 0.5, 0.3, 12, 0.8, size=24, bold=True, color=navy)
    summary = report.get("executive_summary", "")[:600]
    _text_box(s2, summary, 0.5, 1.3, 12, 3.5, size=16, color=navy)

    # KPI boxes
    kpis = report.get("kpis", [])[:4]
    for i, kpi in enumerate(kpis):
        left = 0.5 + i * 3.0
        _text_box(
            s2, str(kpi.get("value", "")), left, 5.2, 2.8, 0.8, size=28, bold=True, color=violet
        )
        _text_box(s2, kpi.get("name", ""), left, 6.0, 2.8, 0.5, size=12, color=navy)
    _emit(2, "Executive Summary")

    # ── Slide 3: Data Quality ─────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    _text_box(s3, "Data Quality Overview", 0.5, 0.3, 12, 0.8, size=24, bold=True, color=navy)

    kpis_all = report.get("kpis", [])
    completeness = next((k for k in kpis_all if "Completeness" in k.get("name", "")), None)
    duplicates = next((k for k in kpis_all if "Duplicate" in k.get("name", "")), None)
    anomaly_count = len(report.get("anomaly_alerts", []))

    quality_text = (
        f"✓  Completeness: {completeness.get('value', 'N/A')}{'%' if completeness else ''}\n"
        f"✓  Duplicate rows: {duplicates.get('value', 'N/A') if duplicates else 0}\n"
        f"{'⚠' if anomaly_count > 0 else '✓'}  Anomalies detected: {anomaly_count}\n"
        f"✓  Has forecasts: {'Yes' if report.get('has_forecasts') else 'No'}"
    )
    _text_box(s3, quality_text, 0.5, 1.5, 12, 4, size=18, color=navy)
    _emit(3, "Data Quality")

    # ── Slides 4-5: Insights ──────────────────────────────────────────────
    insights = report.get("insights", [])
    for batch_start, slide_num, label in [(0, 4, "Insights 1-3"), (3, 5, "Insights 4-5")]:
        si = prs.slides.add_slide(blank_layout)
        _text_box(si, label, 0.5, 0.3, 12, 0.8, size=24, bold=True, color=navy)
        for j, ins in enumerate(insights[batch_start : batch_start + 3]):
            top = 1.3 + j * 1.9
            _text_box(
                si, ins.get("headline", ""), 0.5, top, 12, 0.7, size=15, bold=True, color=violet
            )
            _text_box(
                si, ins.get("explanation", "")[:180], 0.5, top + 0.7, 12, 1.2, size=12, color=navy
            )
        _emit(slide_num, label)

    # ── Slide 6: Recommendations ──────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    _text_box(s6, "Recommendations", 0.5, 0.3, 12, 0.8, size=24, bold=True, color=navy)
    for i, rec in enumerate(report.get("recommendations", [])[:3]):
        top = 1.3 + i * 2.0
        label_text = f"[{rec.get('priority', '?').upper()}] {rec.get('title', '')}"
        _text_box(s6, label_text, 0.5, top, 12, 0.7, size=15, bold=True, color=violet)
        _text_box(s6, rec.get("action", "")[:180], 0.5, top + 0.7, 12, 1.3, size=12, color=navy)
    _emit(6, "Recommendations")

    # ── Slide 7: Next Steps ───────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank_layout)
    _fill_bg(s7, "5B4FE8")
    _text_box(s7, "Next Steps", 0.5, 0.5, 12, 1, size=32, bold=True, color=white)
    next_steps = (
        "1.  Review the top 3 recommendations with your data team\n"
        "2.  Set up automated anomaly monitoring for flagged columns\n"
        "3.  Schedule a follow-up analysis in 30 days to track improvements"
    )
    _text_box(s7, next_steps, 0.5, 1.8, 12, 3.5, size=18, color=white)
    _text_box(
        s7,
        "Questions?  Export the full report as PDF or Excel.",
        0.5,
        5.8,
        12,
        0.8,
        size=14,
        color=white,
    )
    _emit(7, "Next Steps")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()

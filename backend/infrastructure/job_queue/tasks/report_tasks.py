"""Celery tasks for report generation.

Queue: ``reports``
Typical worker flags: ``--concurrency=2`` (disk-bound; WeasyPrint / openpyxl)

Triggered by:
- ``ExportReportUseCase`` when the user clicks a download button
- ``ReportAgent`` at the end of the standard analysis pipeline
  (auto-generates a PDF for every completed analysis run)

Supported formats:
    pdf    — WeasyPrint renders an HTML Jinja2 template
    xlsx   — openpyxl multi-sheet workbook
    pptx   — python-pptx branded slide deck
    json   — raw InsightReport dict (for downstream pipeline integrations)
"""

from __future__ import annotations

import asyncio
import time
from typing import TYPE_CHECKING, Any

import structlog
from backend.infrastructure.job_queue.celery_app import celery_app

if TYPE_CHECKING:
    from celery import Task

logger = structlog.get_logger(__name__)

SUPPORTED_FORMATS = {"pdf", "xlsx", "pptx", "json"}


@celery_app.task(
    bind=True,
    name="reports.generate",
    max_retries=2,
    default_retry_delay=30,
    acks_late=True,
    soft_time_limit=180,
    time_limit=240,
)
def generate_report(
    self: Task,
    dataset_id: str,
    session_id: str,
    format: str,
    report_id: str | None = None,
) -> dict[str, Any]:
    """Generate and upload an export report for a completed InsightReport.

    Args:
        dataset_id:  Source dataset UUID.
        session_id:  Analysis session whose InsightReport is to be exported.
        format:      Output format: ``'pdf'`` | ``'xlsx'`` | ``'pptx'`` | ``'json'``.
        report_id:   Optional InsightReport UUID (resolved from session if None).

    Returns:
        Dict: ``{storage_key, format, download_url, status}``.

    Raises:
        ValueError:   Unsupported format string.
        RuntimeError: InsightReport not found for the given session.
    """
    if format not in SUPPORTED_FORMATS:
        raise ValueError(
            f"Unsupported report format '{format}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_FORMATS))}"
        )

    logger.info(
        "report_generation_start",
        task_id=self.request.id,
        dataset_id=dataset_id,
        format=format,
    )
    start = time.monotonic()

    try:
        result = asyncio.run(
            _generate_report_async(
                task=self,
                dataset_id=dataset_id,
                session_id=session_id,
                format=format,
                report_id=report_id,
            )
        )
        duration = round(time.monotonic() - start, 2)
        logger.info(
            "report_generation_complete",
            dataset_id=dataset_id,
            format=format,
            duration_seconds=duration,
            storage_key=result.get("storage_key"),
        )
        return result

    except Exception as exc:
        logger.error(
            "report_generation_failed",
            dataset_id=dataset_id,
            format=format,
            error=str(exc),
            attempt=self.request.retries + 1,
        )
        if self.request.retries < self.max_retries:
            raise self.retry(exc=exc) from exc
        raise


async def _generate_report_async(
    task: Task,
    dataset_id: str,
    session_id: str,
    format: str,
    report_id: str | None,
) -> dict[str, Any]:
    """Core async report generation logic."""
    from backend.infrastructure.cache.redis_cache_adapter import get_redis_cache
    from backend.infrastructure.storage.s3_storage_adapter import S3StorageAdapter
    from backend.shared.utils.uuid_factory import new_uuid

    cache = get_redis_cache()
    storage = S3StorageAdapter()

    # ── Load InsightReport from cache (fast) or Postgres (fallback) ───────
    cached = await cache.get_json(f"insights:{dataset_id}")
    report_data = cached if isinstance(cached, dict) else None
    if not report_data:
        report_data = await _load_report_from_db(dataset_id, session_id)

    if not report_data:
        raise RuntimeError(
            f"InsightReport not found for dataset '{dataset_id}', "
            f"session '{session_id}'. Run analysis first."
        )

    # ── Generate the report bytes ─────────────────────────────────────────
    report_bytes = await _render_report(report_data, format)
    content_type = _content_type(format)

    # ── Upload to S3 / MinIO ──────────────────────────────────────────────
    export_id = new_uuid()
    storage_key = f"reports/{dataset_id}/{export_id}.{format}"

    import io

    await storage.upload_fileobj(
        io.BytesIO(report_bytes),
        storage_key,
        content_type=content_type,
    )

    # ── Generate presigned download URL ───────────────────────────────────
    download_url = await storage.generate_presigned_download_url(storage_key)

    # ── Notify frontend via Redis pub/sub ─────────────────────────────────
    import json

    await cache.publish(
        f"dataset:{dataset_id}",
        json.dumps(
            {
                "type": "report.ready",
                "dataset_id": dataset_id,
                "format": format,
                "storage_key": storage_key,
                "download_url": download_url,
            }
        ),
    )

    # ── Cache the download URL (1 hour — matches presigned URL TTL) ────────
    await cache.set(f"report:{export_id}:url", download_url, ttl=900)

    return {
        "storage_key": storage_key,
        "format": format,
        "download_url": download_url,
        "export_id": export_id,
        "status": "ready",
    }


async def _load_report_from_db(dataset_id: str, session_id: str) -> dict | None:
    """Fallback: load the InsightReport dict from Postgres."""
    try:
        from backend.infrastructure.persistence.database import get_session
        from backend.infrastructure.persistence.repositories.postgres_insight_repository import (
            PostgresInsightRepository,
        )

        async with get_session() as db_session:
            repo = PostgresInsightRepository(db_session)
            report = await repo.get_by_dataset_id(dataset_id)
            return report.to_dict() if report else None
    except Exception as exc:
        logger.warning("report_db_load_failed", dataset_id=dataset_id, error=str(exc))
        return None


async def _render_report(report_data: dict, format: str) -> bytes:
    """Dispatch to the appropriate renderer and return raw bytes."""
    if format == "json":
        import json

        return json.dumps(report_data, indent=2, default=str).encode("utf-8")

    if format == "xlsx":
        return await _render_xlsx(report_data)

    if format == "pdf":
        return await _render_pdf(report_data)

    if format == "pptx":
        return await _render_pptx(report_data)

    raise ValueError(f"No renderer for format '{format}'")


async def _render_xlsx(report_data: dict) -> bytes:
    """Generate an Excel workbook from the InsightReport."""
    import asyncio

    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _render_xlsx_sync, report_data)


def _render_xlsx_sync(report_data: dict) -> bytes:
    import io

    from openpyxl import Workbook

    wb = Workbook()

    # ── Sheet 1: KPIs ────────────────────────────────────────────────────
    ws = wb.active
    ws.title = "KPIs"
    ws.append(["Name", "Value", "Unit"])
    for kpi in report_data.get("kpis", []):
        val = kpi.get("value") or {}
        ws.append([kpi.get("name", ""), val.get("raw", ""), val.get("unit", "")])

    # ── Sheet 2: Insights ────────────────────────────────────────────────
    ws2 = wb.create_sheet("Insights")
    ws2.append(["Rank", "Headline", "Impact", "Confidence", "Columns"])
    for i, ins in enumerate(report_data.get("insights", []), 1):
        ws2.append(
            [
                i,
                ins.get("headline", ""),
                ins.get("business_impact", ""),
                ins.get("confidence", ""),
                ", ".join(ins.get("source_columns", [])),
            ]
        )

    # ── Sheet 3: Anomalies ───────────────────────────────────────────────
    ws3 = wb.create_sheet("Anomalies")
    ws3.append(["Column", "Type", "Severity", "Description", "Confidence"])
    for a in report_data.get("anomaly_alerts", []):
        ws3.append(
            [
                a.get("column_name", ""),
                a.get("anomaly_type", ""),
                a.get("severity", ""),
                a.get("description", ""),
                a.get("confidence", ""),
            ]
        )

    # ── Sheet 4: Recommendations ─────────────────────────────────────────
    ws4 = wb.create_sheet("Recommendations")
    ws4.append(["Priority", "Title", "Situation", "Action", "Impact"])
    for r in report_data.get("recommendations", []):
        ws4.append(
            [
                r.get("priority", ""),
                r.get("title", ""),
                r.get("situation", ""),
                r.get("action", ""),
                r.get("estimated_impact", ""),
            ]
        )

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


async def _render_pdf(report_data: dict) -> bytes:
    """Generate a PDF report using WeasyPrint (HTML → PDF)."""
    import asyncio

    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _render_pdf_sync, report_data)


def _render_pdf_sync(report_data: dict) -> bytes:
    """Synchronous PDF rendering — runs in a thread pool executor."""
    try:
        from weasyprint import HTML

        html = _build_pdf_html(report_data)
        return HTML(string=html).write_pdf()
    except ImportError:
        # Graceful degradation — return JSON if WeasyPrint not installed
        import json

        return json.dumps(report_data, indent=2, default=str).encode("utf-8")


def _build_pdf_html(report_data: dict) -> str:
    """Build an HTML string for WeasyPrint to render to PDF."""
    insights_html = "".join(
        f"<div class='insight'>"
        f"<h3>{ins.get('headline', '')}</h3>"
        f"<p>{ins.get('explanation', '')}</p>"
        f"<span class='badge badge-{ins.get('business_impact', 'low')}'>"
        f"{ins.get('business_impact', '')} impact</span>"
        f"</div>"
        for ins in report_data.get("insights", [])[:10]
    )
    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset='utf-8'>
<style>
  body {{ font-family: Arial, sans-serif; margin: 40px; color: #1E1E1C; }}
  h1   {{ color: #5B4FE8; }}
  h2   {{ color: #1A1A3E; border-bottom: 2px solid #E2E0D8; padding-bottom: 6px; }}
  .insight {{ margin: 16px 0; padding: 12px; border-left: 4px solid #5B4FE8; background: #F5F4F1; }}
  .badge {{ display: inline-block; padding: 2px 8px; border-radius: 4px;
            font-size: 11px; font-weight: 600; }}
  .badge-high   {{ background: #FEE2E2; color: #DC2626; }}
  .badge-medium {{ background: #FEF3C7; color: #D97706; }}
  .badge-low    {{ background: #F3F4F6; color: #6B7280; }}
</style>
</head>
<body>
<h1>DataPilot Analysis Report</h1>
<h2>Executive Summary</h2>
<p>{report_data.get("executive_summary", "No summary available.")}</p>
<h2>Top Insights</h2>
{insights_html}
</body>
</html>"""


async def _render_pptx(report_data: dict) -> bytes:
    """Generate a PowerPoint slide deck from the InsightReport."""
    import asyncio

    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _render_pptx_sync, report_data)


def _render_pptx_sync(report_data: dict) -> bytes:
    import io

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        violet = RGBColor(0x5B, 0x4F, 0xE8)
        navy = RGBColor(0x1A, 0x1A, 0x3E)

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        # ── Slide 1: Title ───────────────────────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5)).text_frame
        tf.text = "DataPilot Analysis Report"
        tf.paragraphs[0].runs[0].font.size = Pt(32)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = violet

        # ── Slide 2: Executive Summary ───────────────────────────────────
        slide2 = prs.slides.add_slide(blank_layout)
        title2 = slide2.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
        ).text_frame
        title2.text = "Executive Summary"
        title2.paragraphs[0].runs[0].font.bold = True
        title2.paragraphs[0].runs[0].font.color.rgb = navy

        body2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5)).text_frame
        body2.word_wrap = True
        body2.text = report_data.get("executive_summary", "")

        # ── Slides 3-N: Top insights (one per slide) ──────────────────────
        for ins in report_data.get("insights", [])[:5]:
            sld = prs.slides.add_slide(blank_layout)
            title = sld.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
            ).text_frame
            title.text = ins.get("headline", "")
            title.paragraphs[0].runs[0].font.bold = True
            title.paragraphs[0].runs[0].font.color.rgb = violet

            body = sld.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(3)).text_frame
            body.word_wrap = True
            body.text = ins.get("explanation", "")

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    except ImportError:
        import json

        return json.dumps(report_data, indent=2, default=str).encode("utf-8")


def _content_type(format: str) -> str:
    return {
        "pdf": "application/pdf",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "json": "application/json",
    }.get(format, "application/octet-stream")
